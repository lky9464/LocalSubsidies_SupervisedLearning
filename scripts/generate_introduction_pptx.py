"""
docs/project_introduction.pptx 생성 — 도식·흐름도 포함 소개 슬라이드.

사용: python scripts/generate_introduction_pptx.py
"""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

# ── 색상 · 폰트 ─────────────────────────────────────────────
FONT = "맑은 고딕"
FONT_FALLBACK = "Malgun Gothic"

C_PRIMARY = (0x1B, 0x5E, 0x4A)
C_ACCENT = (0x2E, 0x7D, 0x63)
C_LIGHT = (0xE8, 0xF5, 0xF0)
C_WHITE = (0xFF, 0xFF, 0xFF)
C_TEXT = (0x33, 0x33, 0x33)
C_MUTED = (0x66, 0x66, 0x66)
C_WARN = (0xE6, 0x5C, 0x00)
C_DANGER = (0xC6, 0x28, 0x28)
C_OK = (0x2E, 0x7D, 0x32)
C_BLUE = (0x15, 0x65, 0xC0)


def _rgb(t: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*t)


def build_pptx(out_path: Path | None = None) -> Path:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    out_path = out_path or (ROOT / "docs" / "project_introduction.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── helpers ─────────────────────────────────────────────
    def _font(run, size=14, bold=False, color=C_TEXT):
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)

    def _box(
        slide,
        left,
        top,
        w,
        h,
        text,
        *,
        fill=C_LIGHT,
        line=C_ACCENT,
        size=11,
        bold=False,
        align=PP_ALIGN.CENTER,
    ):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        _font(p.runs[0], size=size, bold=bold)
        return sh

    def _arrow(slide, x, y, h=Inches(0.35)):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x, y, x, y + h
        )
        conn.line.color.rgb = _rgb(C_ACCENT)
        conn.line.width = Pt(2.5)

    def _title_slide(title: str, subtitle: str):
        sl = prs.slides.add_slide(blank)
        # 상단 컬러 바
        bar = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(C_PRIMARY)
        bar.line.fill.background()
        tb = sl.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        _font(p.runs[0], size=32, bold=True, color=C_WHITE)
        tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(2.5))
        for i, line in enumerate(subtitle.split("\n")):
            para = tb2.text_frame.paragraphs[0] if i == 0 else tb2.text_frame.add_paragraph()
            para.text = line
            _font(para.runs[0], size=16, color=C_MUTED)
        badge = _box(
            sl,
            Inches(0.8),
            Inches(5.8),
            Inches(4.5),
            Inches(0.55),
            "로컬 전용 · 127.0.0.1 · 외부 클라우드·과금 없음",
            fill=C_LIGHT,
            size=12,
        )
        return sl

    def _section(title: str, subtitle: str = ""):
        sl = prs.slides.add_slide(blank)
        stripe = sl.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _rgb(C_PRIMARY)
        stripe.line.fill.background()
        tb = sl.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        _font(p.runs[0], size=28, bold=True, color=C_PRIMARY)
        if subtitle:
            tb2 = sl.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12), Inches(0.5))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = subtitle
            _font(p2.runs[0], size=14, color=C_MUTED)
        return sl

    def _bullets(sl, items: list[str], top=Inches(1.6), left=Inches(0.6), width=Inches(12)):
        tb = sl.shapes.add_textbox(left, top, width, Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(8)
            _font(p.runs[0], size=15)

    def _flow_vertical(
        sl,
        steps: list[str],
        *,
        x=Inches(4.5),
        y=Inches(1.5),
        w=Inches(4.2),
        h=Inches(0.65),
        gap=Inches(0.38),
    ):
        cy = y
        for i, step in enumerate(steps):
            _box(sl, x, cy, w, h, step, size=10)
            if i < len(steps) - 1:
                _arrow(sl, x + w / 2, cy + h, gap)
            cy += h + gap

    def _table(sl, headers, rows, top=Inches(1.5), col_w=None):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl = sl.shapes.add_table(n_rows, n_cols, Inches(0.5), top, Inches(12.3), Inches(0.4 * n_rows)).table
        if col_w:
            for i, w in enumerate(col_w):
                tbl.columns[i].width = w
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(C_PRIMARY)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _font(r, size=11, bold=True, color=C_WHITE)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(C_LIGHT)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        _font(r, size=10)

    # ════════════════════════════════════════════════════════
    # 1. 표지
    # ════════════════════════════════════════════════════════
    _title_slide(
        "지방보조금 부정수급\n위험도 측정 시스템",
        "프로젝트 소개 · 로컬 웹 서비스\n"
        "2026-07  |  대상: 업무 담당자·관리자 (비전문가 포함)",
    )

    # ════════════════════════════════════════════════════════
    # 2. 목차
    # ════════════════════════════════════════════════════════
    sl = _section("목차")
    _bullets(
        sl,
        [
            "1.  프로젝트 개발 목적",
            "2.  환경 구성 · PC 요구사항 · raw 데이터 관리",
            "3.  웹 서비스 구동 방법",
            "4.  업무 순서 (흐름도 · 예외 케이스)",
            "5.  웹 메뉴별 설명",
            "6.  학습·평가 vs 추론 · 업무 활용 예시",
        ],
        top=Inches(1.4),
    )

    # ════════════════════════════════════════════════════════
    # 3. 개발 목적 — 배경
    # ════════════════════════════════════════════════════════
    sl = _section("1. 프로젝트 개발 목적", "왜 이 시스템이 필요한가")
    _box(sl, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.2),
         "많은\n보조금 사업 건", fill=(0xFF, 0xF3, 0xE0), size=14, bold=True)
    _box(sl, Inches(4.6), Inches(1.5), Inches(3.8), Inches(2.2),
         "부정수급은\n극소수", fill=(0xFF, 0xEB, 0xEE), size=14, bold=True)
    _box(sl, Inches(8.7), Inches(1.5), Inches(3.8), Inches(2.2),
         "→ 고위험 건\n우선 점검 필요", fill=C_LIGHT, size=14, bold=True)
    _arrow(sl, Inches(2.4), Inches(3.75), Inches(0.3))
    _arrow(sl, Inches(6.5), Inches(3.75), Inches(0.3))
    _box(
        sl, Inches(0.5), Inches(4.3), Inches(12), Inches(2.5),
        "과거 적발·등록·확정 데이터(타겟 있음)로 AI 모델을 학습하고,\n"
        "각 사업 건에 위험도 점수(0~1,000)를 부여해\n"
        "점검 우선순위를 객관적으로 정리합니다.",
        fill=C_WHITE, size=14, align=PP_ALIGN.LEFT,
    )

    # ════════════════════════════════════════════════════════
    # 4. 시스템 처리 흐름 (도식)
    # ════════════════════════════════════════════════════════
    sl = _section("시스템이 하는 일", "데이터 → 모델 → 점수 → 점검 우선순위")
    steps = [
        ("raw CSV\n(학습/추론)", C_LIGHT),
        ("전처리\n·학습", C_ACCENT),
        ("위험도\n점수", C_PRIMARY),
        ("점검\n우선순위표", C_OK),
    ]
    x0 = Inches(0.7)
    bw = Inches(2.6)
    gap = Inches(0.55)
    for i, (txt, col) in enumerate(steps):
        fill = col if isinstance(col, tuple) else C_LIGHT
        fc = C_WHITE if col in (C_ACCENT, C_PRIMARY, C_OK) else C_TEXT
        _box(sl, x0 + i * (bw + gap), Inches(2.0), bw, Inches(1.4), txt,
             fill=fill if fill != C_ACCENT else C_ACCENT,
             line=C_PRIMARY, size=13, bold=True)
        if col in (C_ACCENT, C_PRIMARY, C_OK):
            sh = sl.shapes[-1]
            sh.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
        if i < len(steps) - 1:
            ax = x0 + i * (bw + gap) + bw + Inches(0.05)
            arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, Inches(2.45), gap - Inches(0.1), Inches(0.5))
            arr.fill.solid()
            arr.fill.fore_color.rgb = _rgb(C_MUTED)
            arr.line.fill.background()
    _box(sl, Inches(0.7), Inches(4.0), Inches(5.5), Inches(2.8),
         "A. 학습·평가 (raw/)\n· 모델 생성 · 성능 검증\n· 타겟 포착 분포 확인",
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT)
    _box(sl, Inches(6.5), Inches(4.0), Inches(5.5), Inches(2.8),
         "B. 추론 (raw_inference/)\n· 라벨 미지 최신 데이터\n· 점검 우선순위표 · Excel",
         fill=(0xE3, 0xF2, 0xFD), size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 5. 기대 효과
    # ════════════════════════════════════════════════════════
    sl = _section("기대 효과")
    cards = [
        ("점검 효율", "한정 인력으로\n고위험 건 우선"),
        ("객관성", "동일 규칙·모델\n재현 가능"),
        ("보안", "로컬 PC만\n클라우드 없음"),
        ("검증", "Test에서 포착\n품질 확인 후 추론"),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.5) + i * Inches(3.15)
        _box(sl, x, Inches(1.6), Inches(2.9), Inches(0.55), t, fill=C_PRIMARY, size=13, bold=True)
        sh = sl.shapes[-1]
        sh.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
        _box(sl, x, Inches(2.3), Inches(2.9), Inches(1.5), b, size=12)

    # ════════════════════════════════════════════════════════
    # 6. 알고리즘
    # ════════════════════════════════════════════════════════
    sl = _section("사용 알고리즘 (5종 비교)")
    _box(sl, Inches(0.5), Inches(1.5), Inches(5.5), Inches(1.2),
         "★ 주 모델 (primary)\n평가·순위로 선정 · 4×4 1차 기준", fill=C_PRIMARY, size=13, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    _box(sl, Inches(6.3), Inches(1.5), Inches(5.5), Inches(1.2),
         "★ 보조 모델 (aux)\n교차 확인 · 동등급 정렬", fill=C_ACCENT, size=13, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    _box(sl, Inches(0.5), Inches(3.0), Inches(11.3), Inches(1.0),
         "참고: Stacked Ensemble · EasyEnsemble · Gradient Boosting (비교·리포트용)",
         fill=C_LIGHT, size=12)

    # ════════════════════════════════════════════════════════
    # 7. 폴더 구조 (시각)
    # ════════════════════════════════════════════════════════
    sl = _section("2. 환경 구성 — 폴더 구조", "프로그램과 데이터를 분리")
    _box(sl, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.2),
         "LocalSubsidies_SupervisedLearning/\n"
         "  app/        ← 웹 화면\n"
         "  scripts/    ← 파이프라인\n"
         "  configs/    ← 설정\n"
         "  RunWeb.bat  ← 실행\n\n"
         "✓ GitHub 가능 (코드만)",
         fill=C_LIGHT, size=11, align=PP_ALIGN.LEFT, line=C_OK)
    _box(sl, Inches(6.5), Inches(1.4), Inches(6.0), Inches(5.2),
         "LocalSubsidies_ML_Data/  (data_root)\n"
         "  raw/           학습·평가 CSV\n"
         "  raw_inference/ 추론 CSV\n"
         "  processed/     전처리 결과\n"
         "  algorithms/    모델·점수\n"
         "  ops/           운영 DB\n"
         "  runs/          Job 로그\n\n"
         "✗ GitHub·클라oud 금지",
         fill=(0xFF, 0xEB, 0xEE), size=11, align=PP_ALIGN.LEFT, line=C_DANGER)
    _box(sl, Inches(0.4), Inches(6.75), Inches(12.1), Inches(0.55),
         "⚠ 민감 CSV·행단위 점수·모델은 항상 ML_Data(data_root)에만 보관",
         fill=(0xFF, 0xF8, 0xE1), size=12, bold=True)

    # ════════════════════════════════════════════════════════
    # 8. PC 사양
    # ════════════════════════════════════════════════════════
    sl = _section("PC 요구사항", "사양에 따라 작업 능력·소요 시간이 달라집니다")
    _table(
        sl,
        ["항목", "쾌적", "보통", "부족", "업무 영향"],
        [
            ["RAM", "≥32GB", "16~32GB", "<16GB", "05 학습 메모리 부족 가능"],
            ["CPU", "≥8코어", "4~7", "<4", "수 시간~반나절"],
            ["디스크", "≥50GB", "20~50GB", "<20GB", "모델·점수 저장 공간"],
        ],
        top=Inches(1.35),
    )
    for i, (label, col) in enumerate([("쾌적", C_OK), ("보통", C_WARN), ("부족", C_DANGER)]):
        _box(sl, Inches(0.5) + i * Inches(4.1), Inches(3.5), Inches(3.8), Inches(1.0),
             {"쾌적": "5종 일괄 학습 OK", "보통": "수 시간 · 절전 OFF", "부족": "주·보 2종만 권장"}[label],
             fill=col if col != C_OK else C_LIGHT, size=12, bold=True)
    _box(sl, Inches(0.5), Inches(4.8), Inches(12), Inches(0.7),
         "웹 메뉴 「내 PC 사양 체크」에서 본인 PC를 즉시 확인할 수 있습니다.",
         fill=C_LIGHT, size=12)

    # ════════════════════════════════════════════════════════
    # 9. raw 관리
    # ════════════════════════════════════════════════════════
    sl = _section("raw 데이터 관리 (필독)", "개인정보·업무기밀 유출 방지")
    rules = [
        ("위치 분리", "프로젝트 밖\nML_Data"),
        ("GitHub 금지", "raw·점수·모델\n업로드 금지"),
        ("AI 주의", "CSV 샘플\n붙여넣기 금지"),
        ("로컬만", "127.0.0.1\n터널 금지"),
        ("DB 메타만", "파일명·건수\n내용 미저장"),
        ("백업 통제", "암호화·권한\n검토"),
    ]
    for i, (t, b) in enumerate(rules):
        col = i % 3
        row = i // 3
        x = Inches(0.45) + col * Inches(4.15)
        y = Inches(1.45) + row * Inches(2.0)
        _box(sl, x, y, Inches(3.85), Inches(0.5), t, fill=C_PRIMARY, size=11, bold=True)
        sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
        _box(sl, x, y + Inches(0.55), Inches(3.85), Inches(1.2), b, size=11)
    _box(sl, Inches(0.45), Inches(5.7), Inches(12), Inches(1.2),
         "raw/ = 학습·평가 (타겟 있음)  |  raw_inference/ = 추론 (타겟 없어도 됨)\n"
         "스키마: TLS4902R_Layout 동일",
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 10. 최초 설치
    # ════════════════════════════════════════════════════════
    sl = _section("최초 1회 설치")
    _flow_vertical(
        sl,
        [
            "Python 3.10+ 설치",
            "venv + pip install",
            "local.yaml (data_root)",
            "raw/ · raw_inference/ 생성",
            "RunWeb.bat 실행",
        ],
        x=Inches(4.2),
        y=Inches(1.4),
        w=Inches(4.8),
    )

    # ════════════════════════════════════════════════════════
    # 11. 웹 구동
    # ════════════════════════════════════════════════════════
    sl = _section("3. 웹 서비스 구동")
    _flow_vertical(
        sl,
        [
            "① data_root 확인",
            "② RunWeb.bat 더블클릭",
            "③ http://127.0.0.1:8501",
            "④ 상단 Job 배너로 진행 확인",
        ],
        x=Inches(3.8),
        y=Inches(1.35),
        w=Inches(5.5),
        h=Inches(0.7),
    )
    _box(sl, Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.3),
         "종료: RunWeb.bat\n콘솔 창 닫기\n→ 서버·Job 함께 종료",
         fill=(0xFF, 0xEB, 0xEE), size=12, align=PP_ALIGN.LEFT)
    _box(sl, Inches(6.3), Inches(5.5), Inches(6.0), Inches(1.3),
         "주의: 코드 수정 후\nRunWeb.bat 재시작\n브라우저만 새로고침 X",
         fill=(0xFF, 0xF8, 0xE1), size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 12. 전체 업무 흐름
    # ════════════════════════════════════════════════════════
    sl = _section("4. 전반 업무 순서 — 큰 그림")
    _flow_vertical(
        sl,
        ["PC 사양 확인 · raw 배치 · RunWeb.bat"],
        x=Inches(4.0), y=Inches(1.2), w=Inches(5.0), h=Inches(0.55),
    )
    _box(sl, Inches(2.5), Inches(2.5), Inches(8.0), Inches(1.1),
         "A. 모델 학습 및 평가 (라벨 있는 raw/)",
         fill=C_PRIMARY, size=14, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    _arrow(sl, Inches(6.5), Inches(3.65), Inches(0.35))
    _box(sl, Inches(2.5), Inches(4.05), Inches(8.0), Inches(1.1),
         "B. 추론 (라벨 없는 raw_inference/) → 점검 우선순위표",
         fill=C_BLUE, size=14, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    _box(sl, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
         "학습·평가(A)를 먼저 완료한 뒤 추론(B)을 실행합니다.\n"
         "A 없이 B만 실행하면 모델이 없어 추론할 수 없습니다.",
         fill=C_LIGHT, size=13, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 13. 학습 파이프라인 흐름
    # ════════════════════════════════════════════════════════
    sl = _section("A. 학습·평가 — 정상 흐름")
    left_steps = [
        "데이터 등록",
        "01 통합",
        "02 타겟",
        "03 전처리",
        "04 누수",
    ]
    right_steps = [
        "05 학습",
        "06 Feature",
        "07 평가",
        "08 순위",
        "09·10",
    ]
    for i, s in enumerate(left_steps):
        _box(sl, Inches(0.5), Inches(1.3) + i * Inches(0.95), Inches(2.8), Inches(0.7), s, size=10)
        if i < len(left_steps) - 1:
            _arrow(sl, Inches(1.9), Inches(1.3 + i * 0.95 + 0.7), Inches(0.22))
    _box(sl, Inches(3.5), Inches(3.0), Inches(1.5), Inches(0.7), "PASS", fill=C_OK, size=11, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    for i, s in enumerate(right_steps):
        _box(sl, Inches(5.5), Inches(1.3) + i * Inches(0.95), Inches(2.8), Inches(0.7), s, size=10)
        if i < len(right_steps) - 1:
            _arrow(sl, Inches(6.9), Inches(1.3 + i * 0.95 + 0.7), Inches(0.22))
    _box(sl, Inches(8.8), Inches(1.3), Inches(3.8), Inches(4.5),
         "확인 메뉴\n\n· 모델 비교·평가\n· 타겟 포착 분포\n· 대시보드",
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT)
    arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.35), Inches(3.15), Inches(2.0), Inches(0.45))
    arr.fill.solid()
    arr.fill.fore_color.rgb = _rgb(C_ACCENT)
    arr.line.fill.background()

    # ════════════════════════════════════════════════════════
    # 14. 누수 FAIL
    # ════════════════════════════════════════════════════════
    sl = _section("예외: 04 누수점검 FAIL")
    _flow_vertical(
        sl,
        [
            "04 FAIL",
            "의심 피처 확인",
            "제외 선택",
            "「03부터 재개」",
            "04 PASS → 05~",
        ],
        x=Inches(4.0),
        y=Inches(1.3),
        w=Inches(5.0),
        h=Inches(0.62),
    )
    _box(sl, Inches(0.5), Inches(5.8), Inches(12), Inches(0.9),
         "누수 = 모델이 '정답을 미리 아는' 컬럼을 학습 → 점수가 비현실적으로 좋아짐 → 반드시 PASS 후 학습",
         fill=(0xFF, 0xF8, 0xE1), size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 15. 추론 흐름
    # ════════════════════════════════════════════════════════
    sl = _section("B. 추론 — 운영 흐름")
    _flow_vertical(
        sl,
        [
            "모델 존재 (05·06 완료)",
            "raw_inference 업로드",
            "추론 → 실행",
            "결과 확인 · Excel",
            "대시보드 요약",
        ],
        x=Inches(3.9),
        y=Inches(1.25),
        w=Inches(5.2),
        h=Inches(0.65),
    )
    _box(sl, Inches(0.5), Inches(5.5), Inches(12), Inches(1.0),
         "6월 데이터가 늦게 들어오면: CSV 추가 후 11 추론 재실행 (전체 기간 재산출)",
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 16. 케이스별 안내
    # ════════════════════════════════════════════════════════
    sl = _section("케이스별 — 어디서 무엇을 하나")
    _table(
        sl,
        ["상황", "다음 행동", "확인 메뉴"],
        [
            ["처음 시작", "raw 업로드 → 01~", "학습 파이프라인"],
            ["모델 없음", "05까지 실행", "모델 비교·평가"],
            ["성능 확인", "07·08 후", "대시보드"],
            ["점검 기준 검증", "10 후", "타겟 포착 분포"],
            ["2026 미지 데이터", "raw_inference → 11", "추론 → 결과 확인"],
            ["PC 느림", "주·보 2종만", "내 PC 사양 체크"],
        ],
        top=Inches(1.3),
    )

    # ════════════════════════════════════════════════════════
    # 17. 메뉴 구조 (트리)
    # ════════════════════════════════════════════════════════
    sl = _section("5. 웹 메뉴 구조")
    menu_tree = (
        "대시보드\n"
        "데이터 등록\n"
        "▼ 모델 학습 및 평가\n"
        "    ├ 학습 파이프라인\n"
        "    ├ 모델 비교·평가\n"
        "    └ 타겟 포착 분포\n"
        "▼ 추론\n"
        "    ├ 추론 실행\n"
        "    └ 결과 확인\n"
        "Run 이력\n"
        "내 PC 사양 체크\n"
        "사용자 가이드\n"
        "설정"
    )
    _box(sl, Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5), menu_tree,
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT, line=C_PRIMARY)
    _box(sl, Inches(5.3), Inches(1.3), Inches(7.2), Inches(1.2),
         "상단 전역 배너\nJob 이름 · 진행률 % · 취소",
         fill=C_PRIMARY, size=13, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    _box(sl, Inches(5.3), Inches(2.7), Inches(7.2), Inches(4.1),
         "백그라운드 Job: 학습·추론 실행 중\n"
         "다른 메뉴로 이동해도 계속 진행\n"
         "RunWeb.bat 창은 반드시 유지",
         fill=(0xE3, 0xF2, 0xFD), size=12, align=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════
    # 18. 학습 vs 추론
    # ════════════════════════════════════════════════════════
    sl = _section("6. 학습·평가 vs 추론")
    _box(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.3),
         "학습·평가\n\n"
         "· raw/ (타겟 있음)\n"
         "· 모델 만들기\n"
         "· 성능·포착 검증\n"
         "· 타겟 포착 분포\n"
         "· 실제라벨 있음",
         fill=C_LIGHT, size=12, align=PP_ALIGN.LEFT, line=C_PRIMARY)
    _box(sl, Inches(6.6), Inches(1.3), Inches(5.9), Inches(5.3),
         "추론\n\n"
         "· raw_inference/\n"
         "· 점검 대상 선정\n"
         "· 점검 우선순위표\n"
         "· Excel 내보내기\n"
         "· 라벨 없음",
         fill=(0xE3, 0xF2, 0xFD), size=12, align=PP_ALIGN.LEFT, line=C_BLUE)
    arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.15), Inches(3.5), Inches(0.5), Inches(0.6))
    arr.fill.solid()
    arr.fill.fore_color.rgb = _rgb(C_ACCENT)
    arr.line.fill.background()

    # ════════════════════════════════════════════════════════
    # 19. 4×4 우선순위 (그리드)
    # ════════════════════════════════════════════════════════
    sl = _section("점검 우선순위 — 4×4 개념")
    grades = ["A\n(상위1%)", "B\n(1~5%)", "C\n(5~10%)", "D\n(10%↓)"]
    pri_labels = [
        ["1", "2", "3", "4"],
        ["5", "6", "7", "8"],
        ["9", "10", "11", "12"],
        ["13", "14", "15", "16"],
    ]
    gx, gy = Inches(3.2), Inches(1.35)
    cs = Inches(1.35)
    _box(sl, Inches(1.5), Inches(1.35), Inches(1.5), Inches(0.55), "주＼보", fill=C_PRIMARY, size=10, bold=True)
    sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    for j, g in enumerate(grades):
        _box(sl, gx + j * cs, Inches(1.35), cs, Inches(0.55), g, fill=C_ACCENT, size=9, bold=True)
        sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
    for i in range(4):
        _box(sl, Inches(1.5), gy + Inches(0.6) + i * cs, Inches(1.5), cs, grades[i], fill=C_ACCENT, size=9, bold=True)
        sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(C_WHITE)
        for j in range(4):
            pri = int(pri_labels[i][j])
            fill = C_PRIMARY if pri <= 4 else (C_ACCENT if pri <= 8 else C_LIGHT)
            fc = C_WHITE if pri <= 8 else C_TEXT
            _box(sl, gx + j * cs, gy + Inches(0.6) + i * cs, cs, cs,
                 f"우선순위\n{pri}", fill=fill, size=10, bold=(pri <= 4))
            if pri <= 8:
                sl.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(fc)
    _box(sl, Inches(0.5), Inches(6.0), Inches(12), Inches(0.9),
         "주등급이 먼저 · 같은 주등급 안에서 보조등급 순 · 숫자가 작을수록 점검 우선",
         fill=C_LIGHT, size=12)

    # ════════════════════════════════════════════════════════
    # 20–23. 업무 예시
    # ════════════════════════════════════════════════════════
    examples = [
        ("예시 1: 연 1회 모델 갱신",
         ["raw/ 배치", "파이프라인 01~10", "모델 비교·평가", "타겟 포착 확인", "→ 추론에 사용"]),
        ("예시 2: 분기별 선제 점검",
         ["raw_inference/ 1~5월", "추론(주·보)", "Excel 우선순위", "점검 계획 회의", "6월 추가 후 재실행"]),
        ("예시 3: 주·보 불일치",
         ["주A × 보D → 과탐 가능", "사유 확인 후 점검", "주A × 보A → 최우선", "", ""]),
        ("예시 4: PC 사양 부족",
         ["사양 체크 → 부족", "주·보 2종만 선택", "05·11 실행", "4×4에 충분", ""]),
    ]
    for title, steps in examples:
        sl = _section(title)
        for i, s in enumerate([x for x in steps if x]):
            _box(sl, Inches(0.5) + i * Inches(2.5), Inches(1.5), Inches(2.3), Inches(0.85), s, size=10)
            if i < len([x for x in steps if x]) - 1:
                arr = sl.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(0.5) + (i + 1) * Inches(2.5) - Inches(0.22),
                    Inches(1.75),
                    Inches(0.2),
                    Inches(0.35),
                )
                arr.fill.solid()
                arr.fill.fore_color.rgb = _rgb(C_MUTED)
                arr.line.fill.background()

    # ════════════════════════════════════════════════════════
    # 24. 산출물 · 마무리
    # ════════════════════════════════════════════════════════
    sl = _section("산출 파일 위치 (로컬 · GitHub 금지)")
    _table(
        sl,
        ["구분", "경로"],
        [
            ["Test 점수", "data_root/algorithms/{algo}/scores/test/"],
            ["추론 점수", "data_root/algorithms/{algo}/scores/inference/"],
            ["Test 포착", "operations/ops_queue_test.xlsx"],
            ["추론 우선순위", "operations/ops_queue_inference.xlsx"],
            ["집계 리포트", "프로젝트/outputs/reports/"],
        ],
        top=Inches(1.3),
    )

    sl = _section("감사합니다")
    _box(
        sl, Inches(1.5), Inches(2.0), Inches(10), Inches(3.5),
        "지방보조금 부정수급 위험도 측정 시스템\n\n"
        "· 로컬 PC 전용 (127.0.0.1)\n"
        "· 상세: docs/operations_criteria.md\n"
        "· 웹: RunWeb.bat → 사용자 가이드",
        fill=C_LIGHT, size=16, align=PP_ALIGN.CENTER,
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        import subprocess

        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-pptx", "-q"]
        )
    out = build_pptx()
    print(f"[intro-pptx] wrote {out}")


if __name__ == "__main__":
    main()
